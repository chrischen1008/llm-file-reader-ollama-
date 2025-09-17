import streamlit as st
import ollama
import fitz  # PyMuPDF
from opencc import OpenCC
import re
import os
import pandas as pd
from docx import Document
from dotenv import load_dotenv
from pptx import Presentation
import threading
import queue

# ==== 環境設定 ====
load_dotenv()
LLM_MODEL = os.getenv("LLM_MODEL")
print(f"使用的 LLM 模型: {LLM_MODEL}")

# ==== 工具函式 ====
def remove_think_tags(text):
    """移除各種思考標籤和不必要的內容"""
    # 移除思考標籤
    text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)
    text = re.sub(r'<thinking>.*?</thinking>', '', text, flags=re.DOTALL)
    # 移除單獨的開始或結束標籤
    text = re.sub(r'</?think>', '', text, flags=re.IGNORECASE)
    text = re.sub(r'</?thinking>', '', text, flags=re.IGNORECASE)
    
    # 移除破損的標籤（如 </think> 沒有對應的 <think>）
    text = re.sub(r'</?\s*think\s*>', '', text, flags=re.IGNORECASE)
    text = re.sub(r'</?\s*thinking\s*>', '', text, flags=re.IGNORECASE)
    # 移除多餘空白
    text = re.sub(r'\n\s*\n', '\n\n', text)
    return text.strip()

# 初始化簡體轉繁體
cc = OpenCC('s2t')
def enforce_traditional(text):
    return cc.convert(text)

# ==== 文字預處理 ====
def preprocess_text(text, max_length=50000):
    """預處理文字，移除多餘空白和截斷過長內容"""
    # 移除多餘的空白和換行
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\n+', '\n', text)
    
    # 如果文字過長，截斷但保持完整句子
    if len(text) > max_length:
        text = text[:max_length]
        last_period = text.rfind('。')
        if last_period > max_length * 0.8:  # 如果句號位置合理
            text = text[:last_period + 1]
    
    return text.strip()

# ==== 讀取檔案文字 ====
def get_text_from_files(files):
    full_text = ""
    
    for file in files:
        st.write(f"正在處理檔案：{file.name}")
        ext = file.name.split(".")[-1].lower()
        
        try:
            if ext == "pdf":
                pdf_document = fitz.open(stream=file.read(), filetype="pdf")
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    full_text += page.get_text() + "\n"
            
            elif ext in ["xlsx", "xls"]:
                excel_data = pd.read_excel(file, sheet_name=None)
                for sheet_name, sheet_df in excel_data.items():
                    full_text += f"\n=== Sheet: {sheet_name} ===\n"
                    sheet_df = sheet_df.fillna('')
                    sheet_text = sheet_df.astype(str).apply(lambda row: ' '.join(row), axis=1).str.cat(sep="\n")
                    full_text += sheet_text + "\n"
            
            elif ext == "csv":
                csv_data = pd.read_csv(file)
                csv_text = csv_data.astype(str).apply(lambda row: ' '.join(row), axis=1).str.cat(sep="\n")
                full_text += csv_text + "\n"
            
            elif ext == "docx":
                doc = Document(file)
                doc_text = "\n".join([p.text for p in doc.paragraphs])
                full_text += doc_text + "\n"
            elif ext == "pptx":
                prs = Presentation(file)
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    full_text += f"\n=== Slide {i+1} ===\n" + "\n".join(slide_text) + "\n"
            elif ext == "txt":
                text_content = file.read().decode("utf-8", errors="ignore")
                full_text += text_content + "\n"
            else:
                st.warning(f"不支援的檔案格式：{file.name}")
        except Exception as e:
            st.error(f"處理檔案 {file.name} 時發生錯誤：{e}")
    
    return preprocess_text(full_text)

# ==== 分段處理大文檔 ====
def split_text_into_chunks(text, chunk_size=20000, overlap=2000):
    """將長文本分割成重疊的段落"""
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        if end > len(text):
            end = len(text)
        
        # 尋找適當的分割點（句號或段落）
        if end < len(text):
            split_point = text.rfind('。', start, end)
            if split_point == -1:
                split_point = text.rfind('\n', start, end)
            if split_point != -1:
                end = split_point + 1
        
        chunks.append(text[start:end])
        start = end - overlap if end < len(text) else end
    
    return chunks

# ==== 優化的 Ollama 調用 ====
def get_ollama_summary_optimized(text):
    if not text or not text.strip():
        return "沒有可總結的文字。"

    # 更簡潔的 prompt
    system_prompt = """你是文件摘要專家，請用繁體中文輸出系統操作流程重點。不要使用任何思考過程標籤，直接給出最終答案。"""

    user_prompt = f"""
請閱讀我提供的文件（公司ERP系統操作手冊），並依據以下所有規則，將內容整理成一份簡潔、有條理的系統操作流程摘要。

### 輸出規則 ###
1. **核心目標：**
   - 僅專注於「系統操作流程」。

2. **格式要求：**
   - 採用「條列式」呈現。
   - 每頁摘要 10 到 15 個關鍵重點。

3. **內容要求：**
   - 保留關鍵專有名詞與重要數字。
   - 避免重複內容。
   - 忽略與操作流程無關的所有細節與背景資訊。

內容：
{text[:100000]}
"""

    try:
        # 添加優化參數
        response = ollama.chat(
            model=LLM_MODEL,
            messages=[
                {'role': 'system', 'content': system_prompt},
                {'role': 'user', 'content': user_prompt}
            ],
            options={
                'temperature': 0.3,  # 降低隨機性提升速度
            }
        )
        return response['message']['content']
    except Exception as e:
        return f"與 Ollama 溝通時發生錯誤：{e}"

# ==== 流式處理函數 ====
def stream_ollama_summary(text, progress_callback=None):
    """流式生成摘要，實時顯示結果"""
    if not text or not text.strip():
        return "沒有可總結的文字。"

    system_prompt = """你是文件摘要專家，請用繁體中文輸出系統操作流程重點。不要使用任何思考過程標籤，直接給出最終答案。"""
    user_prompt = f"""
請閱讀我提供的文件（公司ERP系統操作手冊），並依據以下所有規則，將內容整理成一份簡潔、有條理的系統操作流程摘要。

### 輸出規則 ###
1. **核心目標：**
   - 僅專注於「系統操作流程」。

2. **格式要求：**
   - 採用「條列式」呈現。
   - 每頁摘要 10 到 15 個關鍵重點。

3. **內容要求：**
   - 保留關鍵專有名詞與重要數字。
   - 避免重複內容。
   - 忽略與操作流程無關的所有細節與背景資訊。

內容：
{text[:100000]}
"""

    try:
        # 使用流式API
        stream = ollama.chat(
            model=LLM_MODEL,
            messages=[
                {'role': 'system', 'content': system_prompt},
                {'role': 'user', 'content': user_prompt}
            ],
            stream=True,
            options={
                'temperature': 0.3,
            }
        )
        
        full_response = ""
        for chunk in stream:
            if 'message' in chunk:
                content = chunk['message']['content']
                full_response += content
                if progress_callback:
                    progress_callback(full_response)
        
        return full_response
    except Exception as e:
        return f"與 Ollama 溝通時發生錯誤：{e}"

# ==== Streamlit 介面 ====
st.set_page_config(page_title="LLM 文件總結器", layout="wide")
st.header("使用 LLM 進行重點整理 (PDF / Excel / CSV / Word / PPTX / TXT)")
st.info("此程式會讀取上傳檔案的內容，然後交由 LLM 生成。")

# 添加處理選項
use_streaming = st.checkbox("使用串流模式（即時顯示結果）", value=True)

uploaded_files = st.file_uploader(
    "請上傳 PDF / Excel / CSV / Word / PPTX / TXT 檔案",
    type=["pdf", "xlsx", "xls", "csv", "docx", "pptx", "txt"],
    accept_multiple_files=True
)

if uploaded_files and st.button("開始總結"):
    with st.spinner("正在讀取檔案內容..."):
        full_text = get_text_from_files(uploaded_files)

    with st.expander("點此查看內容"):
        st.text_area("內容", value=full_text[:5000] + "..." if len(full_text) > 5000 else full_text, height=300)

    if full_text.strip():
        st.subheader("文件總結")
        
        if use_streaming:
            # 串流模式
            placeholder = st.empty()
            
            def update_display(content):
                processed_content = enforce_traditional(remove_think_tags(content))
                placeholder.write(processed_content)
            
            with st.spinner("正在使用 LLM 總結文件（串流模式）..."):
                summary = stream_ollama_summary(full_text, update_display)
        else:
            # 傳統模式
            with st.spinner("正在使用 LLM 總結文件..."):
                summary = get_ollama_summary_optimized(full_text)
            
            summary = enforce_traditional(remove_think_tags(summary))
            st.write(summary)
        
        st.success("總結完成！")
    else:
        st.warning("沒有可總結的文字，請確保檔案內容可被讀取。")